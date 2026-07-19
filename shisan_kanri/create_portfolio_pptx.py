#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import io
import os
import warnings
warnings.filterwarnings('ignore')

# ─── フォント設定 ───────────────────────────────────────────────────────────
import matplotlib.font_manager as fm
available = [f.name for f in fm.fontManager.ttflist]
JP_FONT = next(
    (f for f in ['Hiragino Sans', 'Hiragino Kaku Gothic Pro', 'Hiragino Kaku Gothic ProN',
                 'Yu Gothic', 'Noto Sans CJK JP', 'IPAexGothic']
     if f in available),
    None
)
if JP_FONT:
    plt.rcParams['font.family'] = JP_FONT
else:
    # フォントファイルを直接探す
    for fp in fm.findSystemFonts():
        if any(k in fp for k in ['Hiragino', 'YuGothic', 'NotoSansCJK']):
            fe = fm.FontEntry(fname=fp, name='JP')
            fm.fontManager.ttflist.insert(0, fe)
            plt.rcParams['font.family'] = 'JP'
            JP_FONT = 'JP'
            break

plt.rcParams['axes.unicode_minus'] = False

# ─── カラーパレット ─────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1E, 0x3A, 0x5F)
BLUE   = RGBColor(0x2E, 0x86, 0xC1)
LIGHT  = RGBColor(0xD6, 0xEA, 0xF8)
GREEN  = RGBColor(0x1E, 0x8B, 0x4C)
RED    = RGBColor(0xC0, 0x39, 0x2B)
GRAY   = RGBColor(0x58, 0x58, 0x58)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GOLD   = RGBColor(0xF3, 0x9C, 0x12)
TEAL   = RGBColor(0x17, 0x7B, 0x8B)

PIE_COLORS = ['#2E86C1','#1E8B4C','#F39C12','#8E44AD','#C0392B','#17A589','#2C3E50','#E74C3C']

# ─── データ定義 ─────────────────────────────────────────────────────────────

# 日本株（特定預り）
jp_tokutei = [
    ("ＪＡＣリク", 1804120, 833160),
    ("三井物産",   2037312, 1437264),
    ("三菱商事",    699160, 353892),
    ("武田薬品",   1085480, 258720),
    ("蔵王産業",   1404089, 469362),
    ("ＮＴＴ",    1202130, 186480),
    ("ＫＤＤＩ",   576958, 213166),
    ("オリックス", 1068925, 716898),
    ("三菱ＵＦＪ",  616264, 500832),
    ("Ｅギャランティ", 503700, 99300),
    ("三井住友",    252798, 200109),
    ("ジャックス",  352000, 83700),
    ("東京海上",    129492, 97056),
    ("第一ライフ",  116128, 83104),
    ("ＣＤＳ",     583416, 49950),
    ("バルカー",   128960, 95552),
    ("伊藤忠",     193987, 122902),
    ("三菱ＨＣ",    28864, 13904),
    ("アサンテ",    53770, -16036),
    ("アビスト",    39960, 3516),
    ("沖縄セルラー", 29200, 18816),
    ("ＪＲ九州",   27480, 5992),
]

# NISA成長投資枠（株式）
jp_nisa_growth = [
    ("キヤノン",  431000, -56300),
    ("三菱商事", 454000, 193900),
    ("Ｅギャランティ", 335800, 78600),
    ("ＮＴＴ",   505400, -89600),
]

# 旧NISA
jp_old_nisa = [("信越化学", 731000, 302700)]

# 日本株合計
jp_stock_total = 12934193.5 + 1726200 + 731000   # 15,391,393.5
jp_stock_profit = 5827639.5 + 126600 + 302700     # 6,256,939.5

# 投資信託
jp_funds = [
    ("eMAXIS Slim\nオールカントリー\n（成長投資枠）", 4253472, 1140708, 3112764),
    ("eMAXIS Slim\nオールカントリー\n（つみたて枠）", 3362294, 1062254, 2300040),
]
jp_fund_total   = 4253472 + 3362294    # 7,615,766
jp_fund_profit  = 1140708 + 1062254   # 2,202,962

# 米国株（特定預り）
us_tokutei = [
    ("VTI\nバンガードTSM-ETF",        180, 211.55, 369.99, 10749615, 6510795),
    ("HDV\niシェアーズ高配当ETF",      750, 19.42,  26.99,  3267341,  1646591),
    ("SPYD\nS&P500高配当ETF",          100, 40.32,  47.50,   766697,   317997),
    ("AMBQ\nアンビック マイクロ",        80, 34.98,  90.48,  1168350,   752830),
]
us_tokutei_total  = 15952003
us_tokutei_profit = 9228213

# 米国株（NISA）
us_nisa = [
    ("HDV\niシェアーズ高配当ETF",  440, 21.86, 26.99, 1916840, 494760),
    ("PLTR\nパランティア",           7, 72.00, 128.47, 145154,   68112),
]
us_nisa_total  = 2061994
us_nisa_profit = 562872

# 旧NISA米国株
us_old_nisa = [("SKYT\nスカイウォーター", 450, 5.39, 36.57, 2656243, 2293993)]
us_old_nisa_total  = 2656243
us_old_nisa_profit = 2293993

# 米国株合計
us_stock_total  = us_tokutei_total + us_nisa_total + us_old_nisa_total   # 20,670,240
us_stock_profit = us_tokutei_profit + us_nisa_profit + us_old_nisa_profit  # 12,085,078

# 米国債券
us_bonds = [
    ("米国国債 4.125%\n2027/9満期",  12717.44, 2052721),
    ("米国国債 4.750%\n2053/11満期", 13178.10, 2127077),
    ("米国国債 4.250%\n2035/8満期",  3900.00,   629499),
    ("米国国債 4.000%\n2035/11満期", 14626.28, 2360827),
]
us_bond_total = 7170124

# ポートフォリオ全体
GRAND_TOTAL  = jp_stock_total + jp_fund_total + us_stock_total + us_bond_total
GRAND_PROFIT = jp_stock_profit + jp_fund_profit + us_stock_profit  # 債券は損益なし

# ─── ユーティリティ ─────────────────────────────────────────────────────────

def fig_to_image(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor=fig.get_facecolor())
    buf.seek(0)
    plt.close(fig)
    return buf


def add_slide(prs, layout_idx=6):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


def set_bg(slide, color=RGBColor(0xF8, 0xFB, 0xFF)):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, text, left, top, width, height,
                 font_size=14, bold=False, color=None,
                 align=PP_ALIGN.LEFT, font_name="Hiragino Sans"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_header_bar(slide, title, subtitle=""):
    """上部ヘッダーバー"""
    from pptx.util import Inches, Pt
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0), Inches(10), Inches(1.1)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    tf = bar.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "  " + title
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Hiragino Sans"

    if subtitle:
        add_text_box(slide, subtitle, 0.2, 0.78, 9.6, 0.35,
                     font_size=10, color=RGBColor(0xAA, 0xCC, 0xEE), align=PP_ALIGN.LEFT)


def add_kpi_card(slide, left, top, width, height, label, value, sub="", positive=None):
    rect = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    rect.fill.solid()
    rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)

    add_text_box(slide, label, left+0.1, top+0.05, width-0.2, 0.3,
                 font_size=9, color=GRAY, align=PP_ALIGN.CENTER)
    val_color = GREEN if positive is True else (RED if positive is False else NAVY)
    add_text_box(slide, value, left+0.05, top+0.3, width-0.1, 0.5,
                 font_size=15, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    if sub:
        add_text_box(slide, sub, left+0.05, top+0.75, width-0.1, 0.25,
                     font_size=8, color=GRAY, align=PP_ALIGN.CENTER)


# ─── チャート生成 ────────────────────────────────────────────────────────────

def make_asset_pie():
    labels = ['日本株\n（特定・NISA）', '日本\n投資信託', '米国株\n（特定・NISA）', '米国\n国債']
    sizes  = [jp_stock_total, jp_fund_total, us_stock_total, us_bond_total]
    colors = ['#2E86C1', '#1E8B4C', '#E67E22', '#8E44AD']
    explode = [0.03]*4

    fig, ax = plt.subplots(figsize=(5.5, 4.5), facecolor='#F8FBFF')
    wedges, texts, autotexts = ax.pie(
        sizes, labels=labels, colors=colors, explode=explode,
        autopct='%1.1f%%', startangle=140,
        textprops={'fontsize': 9},
        wedgeprops={'linewidth': 1.5, 'edgecolor': 'white'}
    )
    for at in autotexts:
        at.set_fontsize(8)
        at.set_fontweight('bold')
    ax.set_title('アセット配分', fontsize=13, fontweight='bold', pad=10)
    fig.tight_layout()
    return fig_to_image(fig)


def make_account_bar():
    labels = ['日本株\n特定預り', '日本株\nNISA成長', '日本株\n旧NISA',
              'eMAXIS\n成長投資枠', 'eMAXIS\nつみたて枠',
              '米国株\n特定預り', '米国株\nNISA', '米国株\n旧NISA', '米国\n国債']
    values = [jp_stock_total - 1726200 - 731000,
              1726200, 731000,
              4253472, 3362294,
              us_tokutei_total, us_nisa_total, us_old_nisa_total, us_bond_total]
    colors = ['#2E86C1','#5DADE2','#85C1E9',
              '#1E8B4C','#52BE80',
              '#E67E22','#F0A500','#D4AC0D','#8E44AD']

    fig, ax = plt.subplots(figsize=(9, 3.5), facecolor='#F8FBFF')
    bars = ax.bar(range(len(labels)), [v/1e6 for v in values],
                  color=colors, edgecolor='white', linewidth=1.2, width=0.65)
    ax.set_xticks(range(len(labels)))
    ax.set_xticklabels(labels, fontsize=8)
    ax.set_ylabel('評価額（百万円）', fontsize=9)
    ax.set_title('口座別 評価額', fontsize=12, fontweight='bold')
    ax.grid(axis='y', alpha=0.4)
    ax.spines[['top','right']].set_visible(False)
    for bar, val in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.1,
                f'{val/1e6:.1f}M', ha='center', va='bottom', fontsize=7.5, fontweight='bold')
    fig.tight_layout()
    return fig_to_image(fig)


def make_jp_top10_bar():
    all_stocks = jp_tokutei + jp_nisa_growth + jp_old_nisa
    sorted_stocks = sorted(all_stocks, key=lambda x: x[1], reverse=True)[:12]
    names = [s[0] for s in sorted_stocks]
    vals  = [s[1]/1e4 for s in sorted_stocks]
    profs = [s[2] for s in sorted_stocks]
    colors = ['#2ECC71' if p >= 0 else '#E74C3C' for p in profs]

    fig, ax = plt.subplots(figsize=(8, 4), facecolor='#F8FBFF')
    bars = ax.barh(range(len(names)), vals, color=['#2E86C1']*len(names),
                   edgecolor='white', linewidth=1)
    ax.set_yticks(range(len(names)))
    ax.set_yticklabels(names, fontsize=9)
    ax.set_xlabel('評価額（万円）', fontsize=9)
    ax.set_title('日本株 評価額 上位12銘柄', fontsize=12, fontweight='bold')
    ax.invert_yaxis()
    ax.grid(axis='x', alpha=0.4)
    ax.spines[['top','right']].set_visible(False)
    for i, (bar, val, prof) in enumerate(zip(bars, vals, profs)):
        sign = '+' if prof >= 0 else ''
        pcolor = '#1E8B4C' if prof >= 0 else '#C0392B'
        ax.text(bar.get_width() + 0.5, bar.get_y() + bar.get_height()/2,
                f'{val:.0f}万  {sign}{prof/1e4:.1f}万',
                va='center', fontsize=7.5, color=pcolor)
    ax.set_xlim(0, max(vals)*1.35)
    fig.tight_layout()
    return fig_to_image(fig)


def make_jp_sector_pie():
    sectors = {
        '商社': 2037312 + (699160+454000) + 193987,
        '金融': 616264 + 252798 + 352000 + 1068925 + 28864 + 116128 + 129492 + (503700+335800),
        '通信': (1202130+505400) + 576958 + 29200,
        '化学・素材': 1085480 + 731000,
        '人材・ｻｰﾋﾞｽ': 1804120 + 53770 + 39960 + 27480,
        'その他': 583416 + 128960 + 1404089 + 431000,
    }
    fig, ax = plt.subplots(figsize=(5.5, 4.5), facecolor='#F8FBFF')
    wedges, texts, autotexts = ax.pie(
        list(sectors.values()),
        labels=list(sectors.keys()),
        colors=['#2E86C1','#1E8B4C','#E67E22','#8E44AD','#C0392B','#17A589'],
        autopct='%1.1f%%', startangle=90,
        textprops={'fontsize': 10},
        wedgeprops={'linewidth': 1.5, 'edgecolor': 'white'},
        explode=[0.03]*len(sectors)
    )
    for at in autotexts:
        at.set_fontsize(9)
        at.set_fontweight('bold')
    ax.set_title('日本株 セクター構成', fontsize=13, fontweight='bold', pad=10)
    fig.tight_layout()
    return fig_to_image(fig)


def make_fund_chart():
    fig, axes = plt.subplots(1, 2, figsize=(8, 4), facecolor='#F8FBFF')
    for ax, (name, val, profit, cost) in zip(axes, jp_funds):
        slices = [cost, profit]
        clrs   = ['#2E86C1', '#2ECC71']
        wedges, texts, autotexts = ax.pie(
            slices, colors=clrs,
            autopct='%1.1f%%', startangle=90,
            wedgeprops={'linewidth': 1.5, 'edgecolor': 'white'},
            textprops={'fontsize': 9}
        )
        for at in autotexts:
            at.set_fontsize(9)
            at.set_fontweight('bold')
        short_name = name.replace('\n', ' ')
        ax.set_title(f'{short_name}\n評価額 {val/1e4:.0f}万円', fontsize=9, fontweight='bold')
    patches = [mpatches.Patch(color='#2E86C1', label='取得金額'),
               mpatches.Patch(color='#2ECC71', label='評価益')]
    fig.legend(handles=patches, loc='lower center', ncol=2, fontsize=9)
    fig.suptitle('投資信託（eMAXIS Slim オールカントリー）', fontsize=12, fontweight='bold')
    fig.tight_layout(rect=[0, 0.08, 1, 1])
    return fig_to_image(fig)


def make_us_stock_bar():
    stocks = us_tokutei + us_nisa + us_old_nisa
    names  = [s[0].split('\n')[0] for s in stocks]
    vals   = [s[4]/1e6 for s in stocks]
    profs  = [s[5] for s in stocks]
    colors = ['#E67E22','#E67E22','#E67E22','#E67E22',
              '#F0A500','#F0A500','#D4AC0D']

    fig, ax = plt.subplots(figsize=(8, 4), facecolor='#F8FBFF')
    bars = ax.bar(range(len(names)), vals, color=colors,
                  edgecolor='white', linewidth=1.2, width=0.65)
    ax.set_xticks(range(len(names)))
    ax.set_xticklabels(names, fontsize=10)
    ax.set_ylabel('円換算評価額（百万円）', fontsize=9)
    ax.set_title('米国株 銘柄別 評価額（円換算）', fontsize=12, fontweight='bold')
    ax.grid(axis='y', alpha=0.4)
    ax.spines[['top','right']].set_visible(False)
    for bar, val, prof in zip(bars, vals, profs):
        sign = '+' if prof >= 0 else ''
        pcolor = '#1E8B4C' if prof >= 0 else '#C0392B'
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.1,
                f'{val:.2f}M\n{sign}{prof/1e4:.0f}万',
                ha='center', va='bottom', fontsize=8, fontweight='bold', color=pcolor)
    patches = [mpatches.Patch(color='#E67E22', label='特定預り'),
               mpatches.Patch(color='#F0A500', label='NISA'),
               mpatches.Patch(color='#D4AC0D', label='旧NISA')]
    ax.legend(handles=patches, fontsize=9)
    fig.tight_layout()
    return fig_to_image(fig)


def make_bond_chart():
    labels = ['米国国債\n4.125%\n2027満期', '米国国債\n4.750%\n2053満期',
              '米国国債\n4.250%\n2035満期', '米国国債\n4.000%\n2035満期']
    vals = [2052721, 2127077, 629499, 2360827]
    colors = ['#8E44AD','#9B59B6','#A569BD','#BB8FCE']

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(8.5, 3.8), facecolor='#F8FBFF')

    # 棒グラフ
    bars = ax1.bar(range(4), [v/1e4 for v in vals], color=colors,
                   edgecolor='white', linewidth=1.2, width=0.65)
    ax1.set_xticks(range(4))
    ax1.set_xticklabels(labels, fontsize=7.5)
    ax1.set_ylabel('円換算評価額（万円）', fontsize=9)
    ax1.set_title('銘柄別 評価額', fontsize=11, fontweight='bold')
    ax1.grid(axis='y', alpha=0.4)
    ax1.spines[['top','right']].set_visible(False)
    for bar, val in zip(bars, vals):
        ax1.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 3,
                 f'{val/1e4:.0f}万', ha='center', va='bottom', fontsize=8, fontweight='bold')

    # 構成比パイ
    wedges, texts, autotexts = ax2.pie(
        vals, labels=['2027満期','2053満期','2035(8月)満期','2035(11月)満期'],
        colors=colors, autopct='%1.1f%%', startangle=90,
        textprops={'fontsize': 9},
        wedgeprops={'linewidth': 1.5, 'edgecolor': 'white'},
        explode=[0.02]*4
    )
    for at in autotexts:
        at.set_fontsize(9); at.set_fontweight('bold')
    ax2.set_title('満期別 構成比', fontsize=11, fontweight='bold')

    fig.suptitle('米ドル建て国債 ポートフォリオ', fontsize=12, fontweight='bold')
    fig.tight_layout()
    return fig_to_image(fig)


def make_performance_bar():
    categories = ['日本株\n（特定・NISA）', '日本\n投資信託', '米国株\n（全口座）']
    costs   = [jp_stock_total - jp_stock_profit,
               jp_fund_total  - jp_fund_profit,
               us_stock_total - us_stock_profit]
    profits = [jp_stock_profit, jp_fund_profit, us_stock_profit]

    fig, ax = plt.subplots(figsize=(7, 4), facecolor='#F8FBFF')
    x = np.arange(len(categories))
    w = 0.55
    b1 = ax.bar(x, [c/1e6 for c in costs],   width=w, label='取得金額', color='#2E86C1', edgecolor='white')
    b2 = ax.bar(x, [p/1e6 for p in profits],  width=w, bottom=[c/1e6 for c in costs],
                label='評価益', color='#2ECC71', edgecolor='white')

    ax.set_xticks(x)
    ax.set_xticklabels(categories, fontsize=10)
    ax.set_ylabel('金額（百万円）', fontsize=9)
    ax.set_title('取得金額 vs 評価益', fontsize=12, fontweight='bold')
    ax.grid(axis='y', alpha=0.4)
    ax.spines[['top','right']].set_visible(False)
    ax.legend(fontsize=9)

    for bar, p, c in zip(b2, profits, costs):
        rate = p / c * 100
        ax.text(bar.get_x() + bar.get_width()/2,
                (c + p)/1e6 + 0.2,
                f'+{rate:.0f}%', ha='center', fontsize=9,
                fontweight='bold', color='#1E8B4C')
    fig.tight_layout()
    return fig_to_image(fig)


# ─── スライド作成 ─────────────────────────────────────────────────────────────

def create_presentation():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    # ── スライド 1: 表紙 ────────────────────────────────────────────────────
    slide = add_slide(prs, 6)
    set_bg(slide, NAVY)

    # アクセントライン
    line = slide.shapes.add_shape(1, Inches(0), Inches(3.2), Inches(10), Inches(0.06))
    line.fill.solid(); line.fill.fore_color.rgb = GOLD; line.line.fill.background()

    add_text_box(slide, "保有証券 ポートフォリオ分析", 0.8, 1.5, 8.5, 1.0,
                 font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, "レポート", 0.8, 2.4, 8.5, 0.7,
                 font_size=26, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text_box(slide, "2026年6月19日 時点", 0.8, 3.5, 8.5, 0.5,
                 font_size=14, color=RGBColor(0xAA, 0xCC, 0xEE), align=PP_ALIGN.CENTER)

    # KPI サマリー
    kpis = [
        ("総資産評価額", f"約 {GRAND_TOTAL/1e8:.2f} 億円"),
        ("総評価損益",   f"+{GRAND_PROFIT/1e4:.0f} 万円"),
        ("損益率",       f"+{GRAND_PROFIT/(GRAND_TOTAL-GRAND_PROFIT)*100:.1f}%"),
        ("銘柄数",       "日本22 + 米国7"),
    ]
    for i, (lbl, val) in enumerate(kpis):
        card = slide.shapes.add_shape(1,
            Inches(0.5 + i*2.3), Inches(4.3), Inches(2.1), Inches(1.0))
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0x2A, 0x4A, 0x70)
        card.line.color.rgb = GOLD
        add_text_box(slide, lbl, 0.5 + i*2.3 + 0.05, 4.32, 2.0, 0.3,
                     font_size=8, color=RGBColor(0xAA, 0xCC, 0xEE), align=PP_ALIGN.CENTER)
        add_text_box(slide, val, 0.5 + i*2.3 + 0.05, 4.6, 2.0, 0.5,
                     font_size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    add_text_box(slide, "日本株 ・ 米国株 ・ 投資信託 ・ 米国債券", 0.8, 6.0, 8.5, 0.4,
                 font_size=11, color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.CENTER)

    # ── スライド 2: ポートフォリオ全体概要 ─────────────────────────────────
    slide = add_slide(prs, 6)
    set_bg(slide)
    add_header_bar(slide, "ポートフォリオ全体概要",
                   f"総評価額: {GRAND_TOTAL/1e8:.2f}億円  |  総評価益: +{GRAND_PROFIT/1e4:.0f}万円  |  2026年6月19日")

    # KPIカード
    cards = [
        ("日本株合計",   f"{jp_stock_total/1e4:.0f}万円",  f"損益 +{jp_stock_profit/1e4:.0f}万", True),
        ("日本投資信託", f"{jp_fund_total/1e4:.0f}万円",   f"損益 +{jp_fund_profit/1e4:.0f}万",  True),
        ("米国株合計",   f"{us_stock_total/1e4:.0f}万円",  f"損益 +{us_stock_profit/1e4:.0f}万", True),
        ("米国国債",     f"{us_bond_total/1e4:.0f}万円",   "4本保有",  None),
    ]
    for i, (lbl, val, sub, pos) in enumerate(cards):
        add_kpi_card(slide, 0.2 + i*2.45, 1.2, 2.2, 1.1, lbl, val, sub, pos)

    # チャート2つ
    img1 = make_asset_pie()
    slide.shapes.add_picture(img1, Inches(0.2), Inches(2.4), Inches(4.8), Inches(4.4))
    img2 = make_account_bar()
    slide.shapes.add_picture(img2, Inches(5.0), Inches(2.4), Inches(4.8), Inches(4.4))

    # ── スライド 3: 日本株 評価額ランキング ─────────────────────────────────
    slide = add_slide(prs, 6)
    set_bg(slide)
    add_header_bar(slide, "日本株 保有銘柄 評価額ランキング",
                   f"特定預り 22銘柄 + NISA 4銘柄 + 旧NISA 1銘柄  ｜  評価額合計 {jp_stock_total/1e4:.0f}万円  |  損益 +{jp_stock_profit/1e4:.0f}万円")

    img = make_jp_top10_bar()
    slide.shapes.add_picture(img, Inches(0.3), Inches(1.2), Inches(9.4), Inches(5.7))

    # ── スライド 4: 日本株 セクター分析 ────────────────────────────────────
    slide = add_slide(prs, 6)
    set_bg(slide)
    add_header_bar(slide, "日本株 セクター別分析",
                   "全口座（特定預り・NISA成長投資枠・旧NISA）合計 ｜ 金融・商社・通信が中心")

    img_sec = make_jp_sector_pie()
    slide.shapes.add_picture(img_sec, Inches(0.2), Inches(1.2), Inches(5.0), Inches(5.8))

    # セクター別コメント
    sector_comments = [
        ("金融",     "三菱UFJ・三井住友・オリックス等\nBig3メガバンク含む幅広い金融"),
        ("商社",     "三井物産・三菱商事・伊藤忠\nコア3大商社を網羅"),
        ("通信",     "NTT・KDDI・沖縄セルラー\n高配当安定株"),
        ("化学・素材","武田薬品・信越化学\nグローバル競争力高い銘柄"),
        ("人材・ｻｰﾋﾞｽ","JAC Recruit・アサンテ等\nニッチ高収益企業"),
    ]
    for i, (sec, desc) in enumerate(sector_comments):
        y = 1.3 + i * 1.0
        dot = slide.shapes.add_shape(1, Inches(5.4), Inches(y + 0.1), Inches(0.12), Inches(0.12))
        dot.fill.solid(); dot.fill.fore_color.rgb = BLUE; dot.line.fill.background()
        add_text_box(slide, sec, 5.6, y + 0.02, 1.0, 0.28, font_size=11, bold=True, color=NAVY)
        add_text_box(slide, desc, 5.6, y + 0.25, 3.8, 0.55, font_size=8.5, color=GRAY)

    # ── スライド 5: 日本株 NISA口座 ─────────────────────────────────────────
    slide = add_slide(prs, 6)
    set_bg(slide)
    add_header_bar(slide, "日本株 NISA口座 保有状況",
                   f"成長投資枠 評価額 172.6万円（損益 +12.7万）  |  旧NISA 評価額 73.1万円（損益 +30.3万）")

    headers = ["銘柄", "口座", "保有株数", "取得単価", "現在値", "評価額", "評価損益"]
    rows = [
        ("キヤノン",       "NISA成長", "100", "4,873", "4,310", "431,000", "▲56,300"),
        ("三菱商事",       "NISA成長", "100", "2,601", "4,540", "454,000", "+193,900"),
        ("Eギャランティ",  "NISA成長", "200", "1,286", "1,679", "335,800", "+78,600"),
        ("NTT",            "NISA成長", "3,500", "170",   "144", "505,400", "▲89,600"),
        ("信越化学",       "旧NISA",   "100", "4,283", "7,310", "731,000", "+302,700"),
    ]

    col_w = [2.0, 1.1, 0.95, 1.0, 0.95, 1.1, 1.1]
    col_x = [0.15]
    for w in col_w[:-1]:
        col_x.append(col_x[-1] + w)

    # ヘッダー行
    for j, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
        cell = slide.shapes.add_shape(1, Inches(cx), Inches(1.3), Inches(cw), Inches(0.38))
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY; cell.line.fill.background()
        add_text_box(slide, hdr, cx + 0.05, 1.33, cw - 0.1, 0.32,
                     font_size=9.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for i, row in enumerate(rows):
        bg = RGBColor(0xEE, 0xF5, 0xFF) if i % 2 == 0 else WHITE
        for j, (val, cx, cw) in enumerate(zip(row, col_x, col_w)):
            cell = slide.shapes.add_shape(1, Inches(cx), Inches(1.68 + i*0.45),
                                          Inches(cw), Inches(0.43))
            cell.fill.solid(); cell.fill.fore_color.rgb = bg; cell.line.color.rgb = LIGHT
            color = RED if '▲' in str(val) else (GREEN if '+' in str(val) and j == 6 else None)
            add_text_box(slide, val, cx + 0.05, 1.71 + i*0.45, cw - 0.1, 0.35,
                         font_size=9, color=color or NAVY,
                         bold=(j == 6), align=PP_ALIGN.CENTER)

    # 合計行
    y_sum = 1.68 + len(rows) * 0.45
    totals = ["合計", "", "", "", "", "1,926,200", "+429,300"]
    for j, (val, cx, cw) in enumerate(zip(totals, col_x, col_w)):
        cell = slide.shapes.add_shape(1, Inches(cx), Inches(y_sum), Inches(cw), Inches(0.43))
        cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT; cell.line.color.rgb = BLUE
        add_text_box(slide, val, cx + 0.05, y_sum + 0.03, cw - 0.1, 0.35,
                     font_size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # ── スライド 6: 日本 投資信託 ────────────────────────────────────────────
    slide = add_slide(prs, 6)
    set_bg(slide)
    add_header_bar(slide, "日本 投資信託（eMAXIS Slim）",
                   f"NISA成長投資枠 + つみたて投資枠  ｜  合計評価額 {jp_fund_total/1e4:.0f}万円  ｜  評価益 +{jp_fund_profit/1e4:.0f}万円")

    img_fund = make_fund_chart()
    slide.shapes.add_picture(img_fund, Inches(0.3), Inches(1.2), Inches(7.5), Inches(4.5))

    # サイド情報
    info = [
        ("ファンド名",   "eMAXIS Slim 全世界株式（オール・カントリー）"),
        ("成長投資枠",   "評価額 425.3万円 ｜ 基準価額 37,907円 ｜ +114.1万円"),
        ("つみたて枠",   "評価額 336.2万円 ｜ 基準価額 37,907円 ｜ +106.2万円"),
        ("合計評価益率", f"+{jp_fund_profit/(jp_fund_total-jp_fund_profit)*100:.1f}%"),
        ("分配金",       "再投資型"),
    ]
    for i, (k, v) in enumerate(info):
        add_text_box(slide, k, 7.9, 1.5 + i*0.95, 1.5, 0.3, font_size=8.5, bold=True, color=NAVY)
        add_text_box(slide, v, 7.9, 1.8 + i*0.95, 1.95, 0.5, font_size=8.5, color=GRAY)

    # ── スライド 7: 米国株 ───────────────────────────────────────────────────
    slide = add_slide(prs, 6)
    set_bg(slide)
    add_header_bar(slide, "米国株式 保有状況",
                   f"特定預り 4銘柄 + NISA 2銘柄 + 旧NISA 1銘柄  ｜  円換算合計 {us_stock_total/1e4:.0f}万円  ｜  評価益 +{us_stock_profit/1e4:.0f}万円")

    img_us = make_us_stock_bar()
    slide.shapes.add_picture(img_us, Inches(0.3), Inches(1.25), Inches(9.4), Inches(4.0))

    # 特徴コメント
    us_comments = [
        ("VTI",  "米国全市場インデックスETF。最大保有。円換算+651万円の含み益"),
        ("HDV",  "米国高配当ETF。特定+NISAで計1,190株保有。安定配当収入"),
        ("SPYD", "S&P500高配当ETF。100株保有"),
        ("AMBQ", "アンビック マイクロ。個別株。+75.3万円"),
        ("PLTR", "パランティア（AI/データ分析）。NISA保有。取得比 +78.5%"),
        ("SKYT", "スカイウォーターテクノロジー。旧NISA。取得比 +578%の大化け銘柄"),
    ]
    for i, (ticker, desc) in enumerate(us_comments):
        x_off = 0.3 if i < 3 else 5.1
        y_off = 5.4 + (i % 3) * 0.38
        add_text_box(slide, f"■ {ticker}", x_off, y_off, 1.0, 0.3,
                     font_size=9, bold=True, color=RGBColor(0xE6, 0x7E, 0x22))
        add_text_box(slide, desc, x_off + 0.9, y_off, 3.8, 0.35, font_size=8.5, color=GRAY)

    # ── スライド 8: 米国債券 ────────────────────────────────────────────────
    slide = add_slide(prs, 6)
    set_bg(slide)
    add_header_bar(slide, "米ドル建て国債 保有状況",
                   f"特定預り 4本  ｜  円換算合計 {us_bond_total/1e4:.0f}万円  ｜  為替レート 約161円/ドル")

    img_bond = make_bond_chart()
    slide.shapes.add_picture(img_bond, Inches(0.3), Inches(1.2), Inches(9.4), Inches(4.0))

    # 債券テーブル
    b_headers = ["コード", "銘柄（米国国債）", "利率", "満期", "面額(USD)", "円換算評価額"]
    b_rows = [
        ("MF164", "米国国債", "4.125%", "2027/9/30",  "$12,400",  "205万円"),
        ("MR435", "米国国債", "4.750%", "2053/11/15", "$13,000",  "213万円"),
        ("NA215", "米国国債", "4.250%", "2035/8/15",  "$3,900",   "62.9万円"),
        ("NB428", "米国国債", "4.000%", "2035/11/15", "$14,600",  "236万円"),
    ]
    b_col_w = [0.8, 1.4, 0.7, 1.3, 1.1, 1.2]
    b_col_x = [0.3]
    for w in b_col_w[:-1]:
        b_col_x.append(b_col_x[-1] + w)

    for j, (hdr, cx, cw) in enumerate(zip(b_headers, b_col_x, b_col_w)):
        cell = slide.shapes.add_shape(1, Inches(cx), Inches(5.4), Inches(cw), Inches(0.35))
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x5B, 0x2D, 0x8E)
        cell.line.fill.background()
        add_text_box(slide, hdr, cx+0.03, 5.43, cw-0.06, 0.28,
                     font_size=8.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for i, row in enumerate(b_rows):
        bg = RGBColor(0xF5, 0xEE, 0xFF) if i % 2 == 0 else WHITE
        for j, (val, cx, cw) in enumerate(zip(row, b_col_x, b_col_w)):
            cell = slide.shapes.add_shape(1, Inches(cx), Inches(5.75 + i*0.38),
                                          Inches(cw), Inches(0.36))
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            cell.line.color.rgb = RGBColor(0xDD, 0xCC, 0xEE)
            add_text_box(slide, val, cx+0.03, 5.78 + i*0.38, cw-0.06, 0.28,
                         font_size=9, color=RGBColor(0x4A, 0x1A, 0x7A), align=PP_ALIGN.CENTER)

    # ── スライド 9: 運用成績比較 ────────────────────────────────────────────
    slide = add_slide(prs, 6)
    set_bg(slide)
    add_header_bar(slide, "運用成績まとめ",
                   f"総評価額 {GRAND_TOTAL/1e8:.2f}億円  ｜  総評価益 +{GRAND_PROFIT/1e4:.0f}万円  ｜  総損益率 +{GRAND_PROFIT/(GRAND_TOTAL-GRAND_PROFIT)*100:.1f}%")

    img_perf = make_performance_bar()
    slide.shapes.add_picture(img_perf, Inches(0.3), Inches(1.2), Inches(6.5), Inches(4.5))

    # 総合サマリーカード
    summaries = [
        ("日本株（特定+NISA）", f"{jp_stock_total/1e4:.0f}万",   f"+{jp_stock_profit/1e4:.0f}万",
         f"+{jp_stock_profit/(jp_stock_total-jp_stock_profit)*100:.0f}%"),
        ("日本 投資信託",       f"{jp_fund_total/1e4:.0f}万",    f"+{jp_fund_profit/1e4:.0f}万",
         f"+{jp_fund_profit/(jp_fund_total-jp_fund_profit)*100:.0f}%"),
        ("米国株（全口座）",    f"{us_stock_total/1e4:.0f}万",   f"+{us_stock_profit/1e4:.0f}万",
         f"+{us_stock_profit/(us_stock_total-us_stock_profit)*100:.0f}%"),
        ("米国国債",            f"{us_bond_total/1e4:.0f}万",    "—",   "安定運用"),
    ]
    for i, (label, total, profit, rate) in enumerate(summaries):
        y = 1.3 + i * 1.3
        rect = slide.shapes.add_shape(1, Inches(6.9), Inches(y), Inches(2.95), Inches(1.15))
        rect.fill.solid(); rect.fill.fore_color.rgb = WHITE
        rect.line.color.rgb = BLUE
        add_text_box(slide, label, 7.0, y + 0.05, 2.7, 0.3, font_size=9, bold=True, color=NAVY)
        add_text_box(slide, f"評価額: {total}", 7.0, y + 0.35, 2.7, 0.28, font_size=9, color=GRAY)
        add_text_box(slide, f"損益: {profit}  ({rate})", 7.0, y + 0.63, 2.7, 0.28,
                     font_size=9.5, bold=True, color=GREEN if '+' in profit else NAVY)

    # 総合合計
    rect = slide.shapes.add_shape(1, Inches(6.9), Inches(6.55), Inches(2.95), Inches(0.7))
    rect.fill.solid(); rect.fill.fore_color.rgb = NAVY; rect.line.fill.background()
    add_text_box(slide, f"ポートフォリオ合計:  {GRAND_TOTAL/1e8:.2f}億円",
                 7.0, 6.6, 2.7, 0.3, font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, f"総評価益: +{GRAND_PROFIT/1e4:.0f}万円  (+{GRAND_PROFIT/(GRAND_TOTAL-GRAND_PROFIT)*100:.1f}%)",
                 7.0, 6.9, 2.7, 0.28, font_size=9, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # ── スライド 10: ポートフォリオの特徴と注目銘柄 ──────────────────────────
    slide = add_slide(prs, 6)
    set_bg(slide)
    add_header_bar(slide, "ポートフォリオの特徴と注目銘柄")

    features = [
        ("分散投資",    "#2E86C1",
         "日本株22銘柄 ＋ 米国株7銘柄 ＋ 投信 ＋ 国債の4資産に分散。\n地域・セクター・商品タイプを幅広くカバー。"),
        ("高配当重視",  "#1E8B4C",
         "国内：NTT・KDDI・オリックス・三菱商事など高配当株を中心に保有。\n米国：HDV・SPYDの高配当ETFでインカムゲインを確保。"),
        ("成長株保有",  "#E67E22",
         "SKYT（+578%）・PLTR（+78%）など米国ハイテク成長株をNISA口座で保有。\nJAC（+86%）・三井物産（+239%）なども大幅上昇。"),
        ("NISA活用",    "#8E44AD",
         "成長投資枠・つみたて投資枠・旧NISAをフル活用。\n非課税での利益最大化を実現。"),
        ("債券でリスク分散", "#C0392B",
         "米国国債4本（短期・中期・超長期）でポートフォリオの安定化。\n満期分散（2027〜2053年）でリスク管理。"),
    ]

    for i, (title, color_hex, desc) in enumerate(features):
        row = i // 2
        col = i % 2
        if i == 4:  # 最後は中央に
            xl, yl = 2.3, 1.3 + row * 1.9
        else:
            xl, yl = 0.3 + col * 4.9, 1.3 + row * 1.9

        card = slide.shapes.add_shape(1, Inches(xl), Inches(yl), Inches(4.5), Inches(1.7))
        r, g, b = int(color_hex[1:3], 16), int(color_hex[3:5], 16), int(color_hex[5:7], 16)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(
            min(r + 180, 255), min(g + 180, 255), min(b + 180, 255))
        card.line.color.rgb = RGBColor(r, g, b)

        add_text_box(slide, f"● {title}", xl + 0.15, yl + 0.1, 4.1, 0.38,
                     font_size=12, bold=True, color=RGBColor(r, g, b))
        add_text_box(slide, desc, xl + 0.15, yl + 0.48, 4.2, 1.1,
                     font_size=9, color=GRAY)

    # ─── 保存 ──────────────────────────────────────────────────────────────
    output_path = "/Users/takahiroazuma/Desktop/kabu/claude/portfolio_report.pptx"
    prs.save(output_path)
    print(f"✅ 保存完了: {output_path}")
    return output_path


if __name__ == "__main__":
    path = create_presentation()
    size_kb = os.path.getsize(path) // 1024
    print(f"   ファイルサイズ: {size_kb} KB")
    print(f"   スライド数: 10枚")
